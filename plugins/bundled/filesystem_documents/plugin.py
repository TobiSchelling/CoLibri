#!/usr/bin/env python3
"""Filesystem mixed-document ingestion plugin (quality-first).

Contract:
- stdin: { "config": {...}, "cursor": <json|null> }
- stdout: JSONL document envelopes
- optional cursor line at end: { "type": "cursor", "cursor": {...} }
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime, timezone
import fnmatch
import hashlib
import json
from pathlib import Path
import subprocess
import sys
import tempfile

# Works when executed from this folder in the repository checkout.
sys.path.insert(0, str(Path(__file__).resolve().parents[2] / "python-sdk"))
from colibri_plugin_sdk import build_envelope  # noqa: E402


SUPPORTED_DEFAULT_EXTS = [".md", ".markdown", ".pdf", ".epub", ".docx", ".pptx"]


def utc_iso_from_mtime(path: Path) -> str:
    return datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc).isoformat()


def now_utc_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def stable_source_id(root: Path) -> str:
    resolved = root.expanduser().resolve()
    digest = hashlib.sha256(str(resolved).encode("utf-8")).hexdigest()
    return digest[:12]


def stable_doc_id(plugin_id: str, source_id: str, rel_path: str) -> str:
    return f"{plugin_id}:{source_id}:{rel_path}"


def load_request() -> dict:
    raw = sys.stdin.read().strip()
    if not raw:
        return {"config": {}}
    return json.loads(raw)


def which(tool: str) -> str | None:
    p = subprocess.run(
        ["which", tool], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
    )
    return tool if p.returncode == 0 else None


def run_cmd(args: list[str], *, timeout_secs: int = 600) -> subprocess.CompletedProcess:
    return subprocess.run(
        args,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=timeout_secs,
    )


def convert_with_pandoc(input_path: Path, from_format: str) -> str:
    if which("pandoc") is None:
        raise RuntimeError("pandoc not found (install with: brew install pandoc)")

    p = run_cmd(
        [
            "pandoc",
            "-f",
            from_format,
            "-t",
            "gfm",
            "--wrap=none",
            str(input_path),
        ],
        timeout_secs=1200,
    )
    if p.returncode != 0:
        raise RuntimeError(f"pandoc failed: {p.stderr.strip()}")
    return p.stdout


def convert_pdf_with_docling(input_path: Path) -> str:
    if which("docling") is None:
        raise RuntimeError("docling not found (install with: pipx install docling)")

    with tempfile.TemporaryDirectory(prefix="colibri-docling-") as tmp:
        out_dir = Path(tmp)
        p = run_cmd(
            [
                "docling",
                str(input_path),
                "--to",
                "md",
                "--image-export-mode",
                "placeholder",
                "--output",
                str(out_dir),
            ],
            timeout_secs=3600,
        )
        if p.returncode != 0:
            raise RuntimeError(f"docling failed: {p.stderr.strip()}")

        stem = input_path.stem or "output"
        out_md = out_dir / f"{stem}.md"
        if not out_md.exists():
            raise RuntimeError(f"docling did not produce expected output: {out_md}")
        return out_md.read_text(encoding="utf-8", errors="replace")


def convert_pptx_with_soffice_pdf_docling(input_path: Path) -> str:
    if which("soffice") is None:
        raise RuntimeError(
            "soffice not found (install with: brew install libreoffice)"
        )
    if which("docling") is None:
        raise RuntimeError("docling not found (install with: pipx install docling)")

    with tempfile.TemporaryDirectory(prefix="colibri-soffice-") as tmp:
        out_dir = Path(tmp)
        p = run_cmd(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(out_dir),
                str(input_path),
            ],
            timeout_secs=1800,
        )
        if p.returncode != 0:
            raise RuntimeError(f"soffice failed: {p.stderr.strip()}")

        # LibreOffice typically writes <stem>.pdf into outdir, but be robust.
        stem = input_path.stem
        candidates = sorted(out_dir.glob(f"{stem}*.pdf"))
        if not candidates:
            candidates = sorted(out_dir.glob("*.pdf"))
        if not candidates:
            raise RuntimeError("soffice did not produce a PDF output")
        pdf_path = candidates[0]
        return convert_pdf_with_docling(pdf_path)


def convert_pptx_with_python_pptx(input_path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "python-pptx not available (install with: pipx inject <venv> python-pptx or pip install python-pptx)"
        ) from e

    prs = Presentation(str(input_path))
    out: list[str] = []
    out.append(f"# {input_path.stem}")
    out.append("")

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        if hasattr(slide, "shapes") and slide.shapes.title is not None:
            try:
                title = slide.shapes.title.text.strip()
            except Exception:
                title = ""
        out.append(f"## Slide {idx}" + (f": {title}" if title else ""))
        out.append("")

        # Collect text from all shapes (best-effort).
        lines: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
                continue
            try:
                text = shape.text_frame.text.strip()
            except Exception:
                continue
            if text:
                lines.append(text)

        if lines:
            for t in lines:
                out.append(f"- {t}")
            out.append("")

        # Speaker notes (best-effort).
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        if notes:
            out.append("### Notes")
            out.append("")
            out.append(notes)
            out.append("")

    return "\n".join(out).rstrip() + "\n"


def convert_pptx_with_markitdown(input_path: Path) -> str:
    try:
        from markitdown import MarkItDown  # type: ignore
    except Exception as e:  # pragma: no cover
        raise RuntimeError(
            "markitdown not available (install with: pipx install markitdown)"
        ) from e

    md = MarkItDown()
    result = md.convert(str(input_path))
    text = getattr(result, "text_content", None)
    if not isinstance(text, str) or not text.strip():
        raise RuntimeError("markitdown produced empty output")
    return text.rstrip() + "\n"


def convert_pptx(input_path: Path, backend: str) -> str:
    match backend:
        case "soffice_pdf_docling":
            return convert_pptx_with_soffice_pdf_docling(input_path)
        case "pandoc":
            return convert_with_pandoc(input_path, "pptx")
        case "python_pptx":
            return convert_pptx_with_python_pptx(input_path)
        case "markitdown":
            return convert_pptx_with_markitdown(input_path)
        case other:
            raise RuntimeError(f"Unsupported pptx_backend: {other}")


@dataclass(frozen=True)
class PlantUmlSummary:
    entities: list[str]
    relations: list[str]


def parse_plantuml(text: str) -> PlantUmlSummary:
    # Very lightweight parser: alias mapping + arrow relations.
    alias_to_name: dict[str, str] = {}
    entities: set[str] = set()
    relations: list[str] = []

    def norm(s: str) -> str:
        return s.strip().strip('"').strip("'")

    for raw in text.splitlines():
        line = raw.strip()
        if not line or line.startswith("'"):
            continue

        # alias lines: `"User" as U` or `participant "User" as U`
        if " as " in line:
            parts = line.split(" as ", 1)
            left = norm(parts[0].split()[-1] if parts[0].split() else parts[0])
            right = norm(parts[1].split()[0])
            if left and right:
                alias_to_name[right] = left
                entities.add(left)

        # participant declarations: `participant Foo`, `actor "Bar"`
        tokens = line.split()
        if tokens and tokens[0] in {
            "participant",
            "actor",
            "boundary",
            "control",
            "entity",
            "database",
            "component",
            "interface",
            "class",
            "object",
            "usecase",
        }:
            if len(tokens) >= 2:
                name = norm(tokens[1])
                if name:
                    entities.add(name)

        # arrow relations: A -> B : label
        for arrow in ["<->", "<-->", "-->", "->", "<-", "<--", "..>", ".>", "--|>", "-|>"]:
            if arrow in line:
                left, rest = line.split(arrow, 1)
                left = norm(left.split()[-1] if left.split() else left)
                right_part = rest.strip()
                right = norm(right_part.split()[0] if right_part.split() else "")
                if not left or not right:
                    continue
                left = alias_to_name.get(left, left)
                right = alias_to_name.get(right, right)
                entities.add(left)
                entities.add(right)
                label = ""
                if ":" in right_part:
                    label = right_part.split(":", 1)[1].strip()
                rel = f"{left} {arrow} {right}"
                if label:
                    rel = f"{rel}: {label}"
                relations.append(rel)
                break

    return PlantUmlSummary(
        entities=sorted(entities),
        relations=relations,
    )


def strip_existing_plantuml_summaries(md: str) -> str:
    start = "<!-- colibri:plantuml-summary:start -->"
    end = "<!-- colibri:plantuml-summary:end -->"
    out: list[str] = []
    in_block = False
    for line in md.splitlines():
        if line.strip() == start:
            in_block = True
            continue
        if line.strip() == end:
            in_block = False
            continue
        if not in_block:
            out.append(line)
    return "\n".join(out)


def enrich_plantuml_blocks(md: str) -> str:
    md = strip_existing_plantuml_summaries(md)

    lines = md.splitlines()
    out: list[str] = []
    i = 0
    while i < len(lines):
        line = lines[i]
        out.append(line)

        if line.strip().startswith("```") and line.strip() in ("```plantuml", "```puml"):
            fence = line.strip()
            i += 1
            block: list[str] = []
            while i < len(lines) and lines[i].strip() != "```":
                block.append(lines[i])
                out.append(lines[i])
                i += 1
            if i < len(lines):
                out.append(lines[i])  # closing fence
            summary = parse_plantuml("\n".join(block))
            if summary.entities or summary.relations:
                out.append("")
                out.append("<!-- colibri:plantuml-summary:start -->")
                out.append("[CoLibri PlantUML summary]")
                if summary.entities:
                    out.append(f"Entities: {', '.join(summary.entities)}")
                if summary.relations:
                    out.append("Relations:")
                    for r in summary.relations[:50]:
                        out.append(f"- {r}")
                out.append("<!-- colibri:plantuml-summary:end -->")
                out.append("")

        i += 1

    return "\n".join(out).rstrip() + "\n"


def should_exclude(rel_path: str, exclude_globs: list[str]) -> bool:
    for pat in exclude_globs:
        if fnmatch.fnmatch(rel_path, pat):
            return True
    return False


def discover_files(root: Path, include_exts: list[str], exclude_globs: list[str]) -> list[Path]:
    exts = {e.lower() for e in include_exts}
    out: list[Path] = []
    for p in root.rglob("*"):
        if not p.is_file():
            continue
        rel = p.relative_to(root).as_posix()
        if should_exclude(rel, exclude_globs):
            continue
        if p.suffix.lower() in exts:
            out.append(p)
    return sorted(out)


def default_title_for_path(path: Path, rel: str) -> str:
    name = path.stem.replace("_", " ").replace("-", " ").strip()
    return name or rel


def main() -> int:
    req = load_request()
    cfg = req.get("config", {})
    cursor = req.get("cursor") or {}

    root_path = cfg.get("root_path")
    classification = cfg.get("classification", "internal")
    include_extensions = cfg.get("include_extensions") or SUPPORTED_DEFAULT_EXTS
    exclude_globs = cfg.get("exclude_globs") or []
    plantuml_summaries = cfg.get("plantuml_summaries", True)
    pptx_backend = cfg.get("pptx_backend", "soffice_pdf_docling")
    mode = cfg.get("mode", "snapshot")
    doc_type_by_extension = cfg.get("doc_type_by_extension") or {}

    if not root_path:
        print("Missing config.root_path", file=sys.stderr)
        return 2

    root = Path(root_path).expanduser().resolve()
    if not root.exists():
        print(f"root_path does not exist: {root}", file=sys.stderr)
        return 2

    last_scan_at = cursor.get("last_scan_at")
    last_scan_dt = None
    if isinstance(last_scan_at, str):
        try:
            last_scan_dt = datetime.fromisoformat(last_scan_at.replace("Z", "+00:00"))
        except ValueError:
            last_scan_dt = None

    plugin_id = "filesystem_documents"
    source_id = stable_source_id(root)

    files = discover_files(root, include_extensions, exclude_globs)

    for file_path in files:
        rel = file_path.relative_to(root).as_posix()
        ext = file_path.suffix.lower()
        file_mtime = datetime.fromtimestamp(file_path.stat().st_mtime, tz=timezone.utc)
        if mode == "incremental" and last_scan_dt is not None and file_mtime <= last_scan_dt:
            continue

        try:
            if ext in [".md", ".markdown"]:
                markdown = file_path.read_text(encoding="utf-8", errors="replace")
            elif ext == ".pdf":
                markdown = convert_pdf_with_docling(file_path)
            elif ext == ".epub":
                markdown = convert_with_pandoc(file_path, "epub")
            elif ext == ".docx":
                markdown = convert_with_pandoc(file_path, "docx")
            elif ext == ".pptx":
                markdown = convert_pptx(file_path, pptx_backend)
            else:
                continue
        except Exception as e:
            print(f"Failed converting {rel}: {e}", file=sys.stderr)
            continue

        if plantuml_summaries:
            markdown = enrich_plantuml_blocks(markdown)

        doc_type = doc_type_by_extension.get(ext.lstrip("."), "note")
        title = default_title_for_path(file_path, rel)
        envelope = build_envelope(
            plugin_id=plugin_id,
            connector_instance=str(root),
            external_id=rel,
            doc_id=stable_doc_id(plugin_id, source_id, rel),
            title=title,
            markdown=markdown,
            doc_type=doc_type,
            classification=classification,
            uri=str(file_path),
            source_updated_at=file_mtime.isoformat(),
        )
        print(json.dumps(envelope, ensure_ascii=False))

    # Cursor checkpoint (useful for incremental mode).
    print(json.dumps({"type": "cursor", "cursor": {"last_scan_at": now_utc_iso()}}))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
